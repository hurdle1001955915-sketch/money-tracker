#!/usr/bin/env python3
"""
収支管理アプリ 社内プレゼンテーション PowerPoint生成スクリプト
全9スライド（8コンテンツ + 1デモ）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== Color Palette =====
BG_DARK = RGBColor(0x0A, 0x0A, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD_LIGHT = RGBColor(0x22, 0x22, 0x3A)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)       # Purple
ACCENT2 = RGBColor(0x00, 0xD2, 0xFF)      # Cyan
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_ORANGE = RGBColor(0xFF, 0x91, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
WHITE = RGBColor(0xF0, 0xF0, 0xF0)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
DARK_GRAY = RGBColor(0x60, 0x60, 0x70)

# ===== Font =====
FONT_MAIN = "Hiragino Sans"  # macOS Japanese font
FONT_MONO = "Menlo"

# ===== Slide dimensions (16:9) =====
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, font_name=FONT_MAIN,
                alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, default_size=16,
                           default_color=GRAY, font_name=FONT_MAIN,
                           alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple styled lines.
    lines: list of dicts with keys: text, size, color, bold, spacing_after
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = line_data.get("text", "")
        p.text = text
        p.font.size = Pt(line_data.get("size", default_size))
        p.font.color.rgb = line_data.get("color", default_color)
        p.font.bold = line_data.get("bold", False)
        p.font.name = font_name
        p.alignment = line_data.get("alignment", alignment)
        p.space_after = Pt(line_data.get("spacing_after", 4))
        if "line_spacing" in line_data:
            p.line_spacing = Pt(line_data["line_spacing"])

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                      border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_section_label(slide, left, top, text, color=ACCENT):
    """Add a small section label."""
    add_textbox(slide, left, top, 3, 0.3, text,
                font_size=11, color=color, bold=True,
                font_name=FONT_MONO)


def add_stat_badge(slide, left, top, number, label):
    """Add a stat badge (number + label)."""
    card = add_rounded_rect(slide, left, top, 2.2, 1.1, BG_CARD, ACCENT, Pt(1))
    add_textbox(slide, left + 0.1, top + 0.1, 2.0, 0.6, str(number),
                font_size=28, color=ACCENT, bold=True,
                font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, top + 0.65, 2.0, 0.4, label,
                font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # ===================================================================
    # SLIDE 1: Hook / Title
    # ===================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)

    # Main title
    add_textbox(slide1, 1.5, 1.2, 10.5, 1.5,
                "自分のお金の流れ、\nちゃんと見えてますか？",
                font_size=42, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide1, 1.5, 3.0, 10.5, 0.6,
                "収支管理  --  AIと作った、完全ローカル家計簿アプリ",
                font_size=18, color=ACCENT2, alignment=PP_ALIGN.CENTER)

    # Stats row
    stats = [("24,000+", "行のSwiftコード"), ("62", "ファイル"),
             ("2-3", "週間で開発"), ("1", "人で開発")]
    start_x = 2.0
    for i, (num, label) in enumerate(stats):
        add_stat_badge(slide1, start_x + i * 2.5, 4.2, num, label)

    # Presenter
    add_textbox(slide1, 1.5, 6.2, 10.5, 0.4,
                "2026.03.04  |  井原 翔太郎",
                font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ===================================================================
    # SLIDE 2: Why
    # ===================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)

    add_section_label(slide2, 0.8, 0.5, "MOTIVATION")
    add_textbox(slide2, 0.8, 0.85, 8, 0.7,
                'なぜ「自分で」作ったのか',
                font_size=32, color=WHITE, bold=True)

    # Card 1: Privacy (highlighted)
    card1 = add_rounded_rect(slide2, 0.8, 1.8, 11.7, 1.6, BG_CARD, ACCENT_GREEN, Pt(2))
    add_textbox(slide2, 1.2, 1.9, 1.0, 0.5, "🔒", font_size=28)
    add_multiline_textbox(slide2, 2.2, 1.85, 9.8, 1.4, [
        {"text": "プライバシー最優先 -- データは100%デバイス内に保存", "size": 18, "color": WHITE, "bold": True, "spacing_after": 6},
        {"text": "他人が作ったアプリに自分の金融データを預けたくない。", "size": 13, "color": GRAY, "spacing_after": 2},
        {"text": "iCloud同期は明示的にOFF / アナリティクス・テレメトリなし / バックアップもローカルZIPのみ", "size": 12, "color": ACCENT_GREEN, "spacing_after": 0},
    ])

    # Card 2: Frustration
    card2 = add_rounded_rect(slide2, 0.8, 3.7, 5.6, 1.4, BG_CARD)
    add_textbox(slide2, 1.2, 3.8, 0.8, 0.4, "😤", font_size=24)
    add_multiline_textbox(slide2, 2.0, 3.8, 4.0, 1.2, [
        {"text": "既存アプリへの不満", "size": 16, "color": WHITE, "bold": True, "spacing_after": 4},
        {"text": "機能が多すぎ or 少なすぎ。広告も邪魔。\n自分のCSV（PayPay・りそな等）に非対応。", "size": 12, "color": GRAY, "spacing_after": 0},
    ])

    # Card 3: Curiosity
    card3 = add_rounded_rect(slide2, 6.9, 3.7, 5.6, 1.4, BG_CARD)
    add_textbox(slide2, 7.3, 3.8, 0.8, 0.4, "🤖", font_size=24)
    add_multiline_textbox(slide2, 8.1, 3.8, 4.0, 1.2, [
        {"text": "AI開発への好奇心", "size": 16, "color": WHITE, "bold": True, "spacing_after": 4},
        {"text": "「Claude Codeで本格アプリ、\n 本当に作れるの？」を検証したかった。", "size": 12, "color": GRAY, "spacing_after": 0},
    ])

    # Quote
    quote_bg = add_rounded_rect(slide2, 0.8, 5.5, 11.7, 0.7, BG_CARD_LIGHT)
    add_textbox(slide2, 1.2, 5.55, 11.0, 0.6,
                "「自分の金融データは、自分の端末だけで管理する」",
                font_size=16, color=WHITE, bold=False,
                alignment=PP_ALIGN.CENTER)

    # ===================================================================
    # SLIDE 3: 3 Appeal Points
    # ===================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)

    add_section_label(slide3, 0.8, 0.5, "WHY THIS APP")
    add_textbox(slide3, 0.8, 0.85, 8, 0.7,
                "3つの強み",
                font_size=32, color=WHITE, bold=True)

    # Appeal 1: Local
    c1 = add_rounded_rect(slide3, 0.8, 1.8, 3.7, 4.2, BG_CARD, ACCENT_GREEN, Pt(2))
    add_textbox(slide3, 1.1, 1.95, 0.6, 0.5, "1", font_size=36, color=ACCENT_GREEN, bold=True, font_name=FONT_MONO)
    add_textbox(slide3, 1.7, 2.0, 2.5, 0.4, "🔒", font_size=28)
    add_multiline_textbox(slide3, 1.1, 2.7, 3.1, 3.0, [
        {"text": "完全ローカル", "size": 20, "color": WHITE, "bold": True, "spacing_after": 10},
        {"text": "あなたのデータは、\nあなたのデバイスだけに", "size": 14, "color": ACCENT_GREEN, "spacing_after": 10},
        {"text": "SwiftData でデバイス内に永続保存", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "iCloud同期は明示的にOFF", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "アナリティクス・テレメトリなし", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "APIキーはKeychain保存", "size": 11, "color": GRAY, "spacing_after": 0},
    ])

    # Appeal 2: AI
    c2 = add_rounded_rect(slide3, 4.8, 1.8, 3.7, 4.2, BG_CARD, ACCENT, Pt(2))
    add_textbox(slide3, 5.1, 1.95, 0.6, 0.5, "2", font_size=36, color=ACCENT, bold=True, font_name=FONT_MONO)
    add_textbox(slide3, 5.7, 2.0, 2.5, 0.4, "🧠", font_size=28)
    add_multiline_textbox(slide3, 5.1, 2.7, 3.1, 3.0, [
        {"text": "AIが自動仕分け", "size": 20, "color": WHITE, "bold": True, "spacing_after": 10},
        {"text": "CSVインポート +\nGPT-4o-mini自動分類", "size": 14, "color": ACCENT, "spacing_after": 10},
        {"text": "6 CSVフォーマット自動判定", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "信頼度スコア80%以上で自動確定", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "分類理由を日本語で表示", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "ルール学習 + 手動修正対応", "size": 11, "color": GRAY, "spacing_after": 0},
    ])

    # Appeal 3: Visualization
    c3 = add_rounded_rect(slide3, 8.8, 1.8, 3.7, 4.2, BG_CARD, ACCENT2, Pt(2))
    add_textbox(slide3, 9.1, 1.95, 0.6, 0.5, "3", font_size=36, color=ACCENT2, bold=True, font_name=FONT_MONO)
    add_textbox(slide3, 9.7, 2.0, 2.5, 0.4, "📊", font_size=28)
    add_multiline_textbox(slide3, 9.1, 2.7, 3.1, 3.0, [
        {"text": "7つのグラフで\n見える化", "size": 20, "color": WHITE, "bold": True, "spacing_after": 10},
        {"text": "円グラフから\n予算進捗まで", "size": 14, "color": ACCENT2, "spacing_after": 10},
        {"text": "支出/収入カテゴリ円グラフ", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "年間推移棒グラフ・折れ線", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "予算達成プログレス", "size": 11, "color": GRAY, "spacing_after": 4},
        {"text": "資産ダッシュボード", "size": 11, "color": GRAY, "spacing_after": 0},
    ])

    # ===================================================================
    # SLIDE 4: Main Features
    # ===================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)

    add_section_label(slide4, 0.8, 0.5, "FEATURES")
    add_textbox(slide4, 0.8, 0.85, 8, 0.7,
                "主な機能",
                font_size=32, color=WHITE, bold=True)

    # Tab bar mockup
    tabs = [
        ("📝", "入力"),
        ("📅", "カレンダー"),
        ("📊", "グラフ"),
        ("💰", "資産"),
        ("⚙️", "設定"),
    ]
    tab_width = 2.2
    tab_start = 0.8 + (11.7 - tab_width * 5 - 0.3 * 4) / 2
    for i, (icon, name) in enumerate(tabs):
        x = tab_start + i * (tab_width + 0.3)
        add_rounded_rect(slide4, x, 1.7, tab_width, 0.7, BG_CARD, ACCENT, Pt(1))
        add_textbox(slide4, x, 1.75, tab_width, 0.6,
                    f"{icon}  {name}",
                    font_size=14, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Feature grid (2x3)
    features = [
        ("🧮", "電卓式入力", "直感的な金額入力。よく使うカテゴリ\nが自動で上位表示。"),
        ("📷", "レシートOCR", "カメラで撮影 → Vision AIが\n金額と日付を自動読み取り。"),
        ("📅", "カレンダー", "月間カレンダーで日別収支を一覧。\n月次サマリーも表示。"),
        ("💼", "資産ダッシュボード", "総資産額、口座別ポートフォリオ、\n6ヶ月の資産推移チャート。"),
        ("🏧", "8種類の口座管理", "銀行・クレカ・電子マネー・PayPay\n・Suica・現金・投資・その他"),
        ("📋", "予算 & 固定費", "月間予算のカテゴリ別設定。\n固定費テンプレートで自動入力。"),
    ]

    for i, (icon, title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = 0.8 + col * 4.0
        y = 2.8 + row * 2.0
        add_rounded_rect(slide4, x, y, 3.7, 1.7, BG_CARD)
        add_textbox(slide4, x + 0.2, y + 0.15, 0.5, 0.4, icon, font_size=20)
        add_textbox(slide4, x + 0.7, y + 0.15, 2.8, 0.35, title,
                    font_size=14, color=WHITE, bold=True)
        add_textbox(slide4, x + 0.2, y + 0.6, 3.3, 1.0, desc,
                    font_size=11, color=GRAY)

    # Bottom highlight
    add_rounded_rect(slide4, 0.8, 6.4, 11.7, 0.6, BG_CARD_LIGHT)
    add_textbox(slide4, 1.0, 6.42, 11.3, 0.5,
                "🔐 Face ID / Touch ID ロック   |   💾 ZIPバックアップ & リストア   |   🔍 高度な検索 & フィルタ",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ===================================================================
    # SLIDE 5: Architecture (Engineer-focused)
    # ===================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)

    add_section_label(slide5, 0.8, 0.5, "TECH DEEP DIVE")
    add_textbox(slide5, 0.8, 0.85, 8, 0.7,
                "アーキテクチャ概要",
                font_size=32, color=WHITE, bold=True)

    # Left column: Data Model
    add_rounded_rect(slide5, 0.8, 1.8, 5.5, 3.3, BG_CARD, ACCENT, Pt(1))
    add_textbox(slide5, 1.1, 1.9, 3, 0.35, "Data Model",
                font_size=14, color=ACCENT, bold=True, font_name=FONT_MONO)
    add_multiline_textbox(slide5, 1.1, 2.3, 5.0, 2.6, [
        {"text": "Transaction (29 fields)", "size": 13, "color": WHITE, "bold": True, "spacing_after": 2},
        {"text": "  categoryId, classificationSource,", "size": 11, "color": GRAY, "font_name": FONT_MONO, "spacing_after": 1},
        {"text": "  classificationConfidence, fingerprintKey...", "size": 11, "color": GRAY, "font_name": FONT_MONO, "spacing_after": 8},
        {"text": "CategoryGroup → CategoryItem (階層構造)", "size": 12, "color": WHITE, "spacing_after": 4},
        {"text": "Account (8 types) → AccountStore", "size": 12, "color": WHITE, "spacing_after": 4},
        {"text": "Budget / FixedCostTemplate / ClassificationRule", "size": 12, "color": WHITE, "spacing_after": 4},
        {"text": "BackupPayload (v3, 後方互換デコーダ)", "size": 12, "color": WHITE, "spacing_after": 0},
    ], font_name=FONT_MAIN)

    # Right column: Architecture
    add_rounded_rect(slide5, 6.6, 1.8, 5.9, 3.3, BG_CARD, ACCENT2, Pt(1))
    add_textbox(slide5, 6.9, 1.9, 3, 0.35, "Architecture",
                font_size=14, color=ACCENT2, bold=True, font_name=FONT_MONO)
    add_multiline_textbox(slide5, 6.9, 2.3, 5.3, 2.6, [
        {"text": "SwiftUI Views", "size": 13, "color": WHITE, "bold": True, "spacing_after": 2},
        {"text": "    ↓  @EnvironmentObject", "size": 11, "color": ACCENT2, "spacing_after": 2},
        {"text": "DataStore.shared (Singleton)", "size": 13, "color": WHITE, "bold": True, "spacing_after": 2},
        {"text": "    ↓  SwiftData ModelContainer", "size": 11, "color": ACCENT2, "spacing_after": 2},
        {"text": "Local SQLite Storage", "size": 13, "color": ACCENT_GREEN, "bold": True, "spacing_after": 8},
        {"text": "KeychainStore → OpenAI API (opt-in)", "size": 12, "color": WHITE, "spacing_after": 4},
        {"text": "※ 金融データの外部送信なし", "size": 12, "color": ACCENT_GREEN, "bold": True, "spacing_after": 0},
    ], font_name=FONT_MAIN)

    # Bottom stats row
    tech_stats = [
        ("62", "Swiftファイル"),
        ("24,000+", "行のコード"),
        ("5", "テストファイル"),
        ("6", "CSVフォーマット"),
        ("5+", "文字コード対応"),
    ]
    for i, (num, label) in enumerate(tech_stats):
        x = 0.8 + i * 2.45
        add_rounded_rect(slide5, x, 5.4, 2.15, 0.9, BG_CARD)
        add_textbox(slide5, x, 5.42, 2.15, 0.45, num,
                    font_size=20, color=ACCENT, bold=True,
                    font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
        add_textbox(slide5, x, 5.85, 2.15, 0.4, label,
                    font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Tech keywords
    add_textbox(slide5, 0.8, 6.5, 11.7, 0.4,
                "SwiftUI  /  SwiftData  /  Vision  /  LocalAuthentication  /  Security (Keychain)  /  Swift Charts",
                font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                font_name=FONT_MONO)

    # ===================================================================
    # SLIDE 6: Vibe Coding
    # ===================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)

    add_section_label(slide6, 0.8, 0.5, "VIBE CODING")
    add_textbox(slide6, 0.8, 0.85, 10, 0.7,
                "AI駆動開発 -- どうやって作ったのか",
                font_size=30, color=WHITE, bold=True)

    # Tool cards: Claude Code + Google Antigravity
    # Claude Code (primary)
    add_rounded_rect(slide6, 0.8, 1.75, 5.5, 2.0, BG_CARD, ACCENT, Pt(2))
    add_multiline_textbox(slide6, 1.1, 1.85, 5.0, 1.8, [
        {"text": "Claude Code", "size": 20, "color": WHITE, "bold": True, "spacing_after": 2},
        {"text": "メイン開発パートナー（90%）", "size": 12, "color": ACCENT, "bold": True, "spacing_after": 6},
        {"text": "コード生成・アーキテクチャ設計・テスト作成\nリファクタリング・デバッグ → 24,000行の大半を生成", "size": 12, "color": GRAY, "spacing_after": 0},
    ])

    # Google Antigravity
    add_rounded_rect(slide6, 6.7, 1.75, 5.8, 2.0, BG_CARD)
    add_multiline_textbox(slide6, 7.0, 1.85, 5.2, 1.8, [
        {"text": "Google Antigravity", "size": 20, "color": WHITE, "bold": True, "spacing_after": 2},
        {"text": "設計・監査パートナー", "size": 12, "color": ACCENT2, "bold": True, "spacing_after": 6},
        {"text": "設計支援・コード品質監査・ドキュメント生成\nUI方向性レビュー・プレゼン構成策定", "size": 12, "color": GRAY, "spacing_after": 0},
    ])

    # Cursor comparison
    add_rounded_rect(slide6, 0.8, 4.0, 11.7, 1.8, BG_CARD_LIGHT, ACCENT_ORANGE, Pt(1))
    add_textbox(slide6, 1.1, 4.1, 5, 0.3, "Cursorとの比較",
                font_size=14, color=ACCENT_ORANGE, bold=True)

    # Comparison table
    comparisons = [
        ("Cursor", "コード補完中心。エディタ内でリアルタイムに支援。行単位〜関数単位。", GRAY),
        ("Claude Code", "機能単位で自然言語指示 → 複数ファイルに跨る実装を一括生成。", ACCENT),
        ("Google Antigravity", "プロジェクト全体の分析・監査・計画策定。俯瞰的な支援。", ACCENT2),
    ]
    for i, (tool, desc, color) in enumerate(comparisons):
        y = 4.5 + i * 0.4
        add_textbox(slide6, 1.3, y, 2.5, 0.35, tool,
                    font_size=12, color=color, bold=True, font_name=FONT_MONO)
        add_textbox(slide6, 3.8, y, 8.3, 0.35, desc,
                    font_size=12, color=GRAY)

    # Development flow
    add_textbox(slide6, 0.8, 6.0, 2, 0.3, "開発フロー",
                font_size=12, color=ACCENT, bold=True)

    flow_steps = ["💬 日本語で要件", "→", "🤖 AI生成", "→",
                  "🔨 Xcodeビルド", "→", "🔄 エラーFB", "→", "✅ 完成"]
    x = 0.8
    for step in flow_steps:
        if step == "→":
            add_textbox(slide6, x, 6.35, 0.4, 0.4, "→",
                        font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)
            x += 0.4
        else:
            w = 2.0
            add_rounded_rect(slide6, x, 6.3, w, 0.5, BG_CARD)
            add_textbox(slide6, x + 0.05, 6.33, w - 0.1, 0.45, step,
                        font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
            x += w + 0.15

    # ===================================================================
    # SLIDE 7: Future
    # ===================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)

    add_section_label(slide7, 0.8, 0.5, "ROADMAP")
    add_textbox(slide7, 0.8, 0.85, 8, 0.7,
                "今後の展望",
                font_size=32, color=WHITE, bold=True)

    roadmap_items = [
        (ACCENT_ORANGE, "銀行 / クレカ API連携",
         "1円単位の正確な収支把握を自動化。CSV手動インポートからの脱却。", "最重要目標"),
        (ACCENT_GREEN, "iCloud同期の有効化",
         "CloudKitインフラは構築済み。テスト後にフラグONで稼働開始。", "インフラ構築済み"),
        (ACCENT_GREEN, "ウィジェット対応",
         "ホーム画面で今月の収支をさっと確認。WidgetViewsファイル準備済み。", "ファイル準備済み"),
        (ACCENT, "AI支出分析 & 節約アドバイス",
         "支出パターンの予測。「先月より外食が増えています」等の自動通知。", "構想中"),
        (ACCENT, "App Store公開",
         "社内ベータテスト → TestFlight → 一般公開を段階的に検討。", "検討中"),
    ]

    for i, (dot_color, title, desc, status) in enumerate(roadmap_items):
        y = 1.7 + i * 0.95
        # Dot
        dot = slide7.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1.0), Inches(y + 0.15),
                                       Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()

        # Line (except last)
        if i < len(roadmap_items) - 1:
            line_shape = slide7.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.08), Inches(y + 0.35),
                Inches(0.04), Inches(0.6)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(0x30, 0x30, 0x45)
            line_shape.line.fill.background()

        # Content card
        add_rounded_rect(slide7, 1.5, y, 10.5, 0.75, BG_CARD)
        add_textbox(slide7, 1.8, y + 0.05, 6, 0.3, title,
                    font_size=15, color=WHITE, bold=True)
        add_textbox(slide7, 1.8, y + 0.38, 7.5, 0.3, desc,
                    font_size=11, color=GRAY)
        add_textbox(slide7, 9.5, y + 0.1, 2.3, 0.3, status,
                    font_size=10, color=dot_color, bold=True,
                    alignment=PP_ALIGN.RIGHT)

    # ===================================================================
    # SLIDE 8: Summary
    # ===================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)

    add_section_label(slide8, 0.8, 0.5, "TAKEAWAY")
    add_textbox(slide8, 0.8, 0.85, 8, 0.7,
                "まとめ",
                font_size=32, color=WHITE, bold=True)

    # 3 appeal points callback (mini badges)
    badges = [
        ("🔒 完全ローカル", ACCENT_GREEN),
        ("🧠 AI自動仕分け", ACCENT),
        ("📊 7つのグラフ", ACCENT2),
    ]
    for i, (text, color) in enumerate(badges):
        x = 1.5 + i * 3.8
        add_rounded_rect(slide8, x, 1.7, 3.4, 0.6, BG_CARD, color, Pt(1))
        add_textbox(slide8, x + 0.1, 1.75, 3.2, 0.5, text,
                    font_size=14, color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Main message
    add_textbox(slide8, 1.0, 2.8, 11.3, 1.2,
                "AIツールの進化で、\nアイデアと基礎知識があれば\n本格アプリを作れる時代",
                font_size=30, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # CTA
    add_rounded_rect(slide8, 2.5, 4.5, 8.3, 0.8, BG_CARD, ACCENT, Pt(2))
    add_textbox(slide8, 2.8, 4.55, 7.7, 0.7,
                "Cursorの次のステップとして、Claude Codeを試してみてください",
                font_size=16, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Stats footer
    add_textbox(slide8, 1.5, 5.7, 10.3, 0.4,
                "24,000+ lines  /  62 files  /  2-3 weeks  /  1 person  /  0 cloud dependencies",
                font_size=13, color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)

    # Transition text
    add_textbox(slide8, 1.5, 6.5, 10.3, 0.4,
                "それでは、実際にアプリを動かしてお見せします。",
                font_size=16, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # ===================================================================
    # SLIDE 9: Live Demo
    # ===================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9)

    add_textbox(slide9, 1.5, 1.5, 10.3, 1.5,
                "Live Demo",
                font_size=60, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide9, 1.5, 3.2, 10.3, 0.6,
                "iPhone画面を共有します",
                font_size=20, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # Demo flow items
    demo_items = [
        "01  Face ID解除",
        "02  取引入力",
        "03  カレンダー",
        "04  グラフ 7種",
        "05  資産ダッシュボード",
        "06  CSV AI分類",
    ]
    for i, item in enumerate(demo_items):
        col = i % 3
        row = i // 3
        x = 2.0 + col * 3.3
        y = 4.3 + row * 0.8
        add_rounded_rect(slide9, x, y, 2.8, 0.55, BG_CARD, ACCENT, Pt(1))
        add_textbox(slide9, x + 0.1, y + 0.08, 2.6, 0.4, item,
                    font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER,
                    font_name=FONT_MONO)

    add_textbox(slide9, 1.5, 6.3, 10.3, 0.4,
                "ご質問はデモ後にお願いします",
                font_size=13, color=ACCENT2,
                alignment=PP_ALIGN.CENTER)

    # ===================================================================
    # Save
    # ===================================================================
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "slides.pptx")
    prs.save(output_path)
    print(f"PowerPoint saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
